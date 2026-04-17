import os
import io
import streamlit as st
from dotenv import load_dotenv

from modules.ppt_generator import generate_pptx
from modules.slack_client import fetch_channel_messages, search_messages_by_keyword
from modules.web_scraper import scrape_website, search_industry_news

load_dotenv(dotenv_path=os.path.join(os.path.dirname(os.path.abspath(__file__)), ".env"))
load_dotenv(dotenv_path=r"C:\Users\GA\proposal-ppt\.env", override=False)

def _get_secret(key: str, default: str = "") -> str:
    try:
        val = st.secrets.get(key, "")
        if val:
            return val
    except Exception:
        pass
    return os.getenv(key, default)

_ENV_API_KEY = _get_secret("ANTHROPIC_API_KEY")

st.set_page_config(
    page_title="제안서 자동 생성기",
    page_icon="📊",
    layout="wide",
)

st.markdown("""
<style>
.block-container { padding-top: 2rem; }
.source-badge {
    display:inline-block; padding:2px 10px; border-radius:999px;
    background:#DBEAFE; color:#1D4ED8; font-size:0.8rem; margin:2px;
}
.rfp-box {
    background:#FFF7ED; border-left:4px solid #F97316;
    padding:1rem; border-radius:6px; margin-bottom:1rem;
}
</style>
""", unsafe_allow_html=True)

st.title("📊 제안서 자동 생성기")
st.caption("과업지시서와 자료를 넣으면 AI가 논리적인 분석결과를 만들어줍니다")

# ── 파일 읽기 유틸 ────────────────────────────────────────────
def _read_file_safe(f) -> str:
    name = f.name.lower()
    try:
        f.seek(0)
        data = f.read()
        if not data:
            return ""
        if name.endswith(".txt") or name.endswith(".md"):
            return data.decode("utf-8", errors="ignore")
        elif name.endswith(".pdf"):
            import PyPDF2
            reader = PyPDF2.PdfReader(io.BytesIO(data))
            text = "\n".join(p.extract_text() or "" for p in reader.pages)
            if not text.strip():
                return f"[{f.name}] PDF 텍스트 추출 불가 — 스캔 PDF인 경우 텍스트가 없습니다"
            return text
        elif name.endswith(".docx"):
            import docx2txt
            return docx2txt.process(io.BytesIO(data))
        elif name.endswith(".xlsx") or name.endswith(".xls"):
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
            rows = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    line = " | ".join(str(c) for c in row if c is not None)
                    if line.strip():
                        rows.append(line)
            return "\n".join(rows)
        elif name.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            return "\n".join(
                shape.text.strip()
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            )
        elif name.endswith(".hwpx"):
            import zipfile, xml.etree.ElementTree as ET
            with zipfile.ZipFile(io.BytesIO(data)) as z:
                texts = []
                for zname in z.namelist():
                    if "section" in zname.lower() and zname.endswith(".xml"):
                        root = ET.fromstring(z.read(zname))
                        texts += [t for t in root.itertext() if t.strip()]
            return "\n".join(texts)
        elif name.endswith(".hwp"):
            return ""
        else:
            return ""
    except Exception as e:
        st.warning(f"⚠️ {f.name} 읽기 실패: {e}")
        return ""


# ── 상단 핵심 입력 ────────────────────────────────────────────
col_name, col_api = st.columns([2, 2])
with col_name:
    customer_name = st.text_input("🏢 고객사명", placeholder="예) 서울시50플러스재단")
with col_api:
    anthropic_key = st.text_input(
        "🔑 Anthropic API Key",
        type="password",
        placeholder="없으면 비워두세요 (템플릿 모드로 생성)",
    ) or _ENV_API_KEY

if _ENV_API_KEY:
    st.success("🔑 API Key 환경변수 설정됨 — 자동 사용 중")
elif not anthropic_key:
    st.info("💡 API Key 없이도 생성 가능합니다 — AI 분석 없이 입력 내용 기반 기본 구조로 만들어집니다")

st.divider()

# ══════════════════════════════════════════════════════════════
# ★ 과업지시서 섹션 (최우선 입력)
# ══════════════════════════════════════════════════════════════
st.markdown("### 📋 과업지시서 / RFP")
st.caption("고객사로부터 받은 제안요청서, 사업공고문, 과업지시서를 입력하세요 — 가장 중요한 입력입니다")

rfp_col1, rfp_col2 = st.columns([3, 2])

with rfp_col1:
    rfp_text = st.text_area(
        "과업지시서 직접 입력 또는 붙여넣기",
        height=220,
        placeholder=(
            "예)\n"
            "사업명: 2026년 중장년 AI 활용 디지털 역량교육 운영\n"
            "발주기관: 서울시50플러스재단\n"
            "사업기간: 2026.03 ~ 2026.11\n"
            "교육인원: 1,700명 (초급 500명, 중급 800명, 고급 400명)\n"
            "핵심과업:\n"
            "- AI 활용 디지털 역량 교육과정 설계 및 운영\n"
            "- 수준별(초급/중급/고급) 교육과정 개발\n"
            "- 교육 효과성 측정 및 성과 보고\n"
            "평가기준: 교육 만족도 90% 이상, 수료율 85% 이상"
        ),
        label_visibility="collapsed",
    )

with rfp_col2:
    rfp_file = st.file_uploader(
        "과업지시서 파일 업로드",
        type=["txt", "md", "pdf", "docx", "xlsx", "pptx", "hwpx"],
        help="PDF, DOCX, TXT 등 과업지시서 파일을 올려주세요",
        key="rfp_uploader",
    )
    if rfp_file:
        rfp_names = [rfp_file.name]
        if st.session_state.get("rfp_file_names") != rfp_names:
            with st.spinner("과업지시서 읽는 중..."):
                st.session_state["rfp_file_text"] = _read_file_safe(rfp_file)
                st.session_state["rfp_file_names"] = rfp_names
        st.success(f"✅ {rfp_file.name} 로드 완료")
    elif not rfp_file and st.session_state.get("rfp_file_names"):
        st.session_state["rfp_file_text"] = ""
        st.session_state["rfp_file_names"] = []

rfp_file_text = st.session_state.get("rfp_file_text", "")
rfp_combined = "\n\n".join(filter(None, [rfp_text, rfp_file_text]))

if rfp_combined:
    st.success(f"✅ 과업지시서 입력 완료 ({len(rfp_combined)}자) — 이 내용이 제안서 구조의 뼈대가 됩니다")


# ── 추가 자료 입력 탭 ─────────────────────────────────────────
st.markdown("### 📂 추가 참고 자료")
tab1, tab2, tab3 = st.tabs(["📝 메모 / 파일", "💬 Slack", "🌐 웹 데이터"])

with tab1:
    free_text = st.text_area(
        "참고 메모",
        height=130,
        placeholder=(
            "미팅에서 나온 이야기, 고객 Pain Point, 특별히 강조할 내용 등\n\n"
            "예)\n"
            "- 담당자가 특히 강조한 포인트: 수료율 관리, 현업 적용 가능성\n"
            "- 경쟁사 대비 우리 강점: AI 챔피언 인증 프로그램 보유"
        ),
    )

    uploaded_files = st.file_uploader(
        "추가 파일 (회의록, 고객사 자료 등)",
        type=["txt", "md", "pdf", "docx", "xlsx", "pptx", "hwpx"],
        accept_multiple_files=True,
        help="지원: TXT, MD, PDF, DOCX, XLSX, PPTX, HWPX",
    )

    file_text = ""
    if uploaded_files:
        file_names = [f.name for f in uploaded_files]
        if st.session_state.get("uploaded_file_names") != file_names:
            with st.spinner("파일 읽는 중..."):
                file_text = ""
                for f in uploaded_files:
                    content = _read_file_safe(f)
                    file_text += f"\n\n[파일: {f.name}]\n{content}"
                st.session_state["file_text"] = file_text
                st.session_state["uploaded_file_names"] = file_names
        st.success(f"✅ {len(uploaded_files)}개 파일 로드 완료")
    elif not uploaded_files and st.session_state.get("uploaded_file_names"):
        st.session_state["file_text"] = ""
        st.session_state["uploaded_file_names"] = []

    file_text = st.session_state.get("file_text", "")

with tab2:
    slack_token = st.text_input(
        "Slack Bot Token",
        value=os.getenv("SLACK_BOT_TOKEN", ""),
        type="password",
        placeholder="xoxb-... (없으면 건너뛰기)",
    )
    if slack_token:
        s1, s2 = st.columns(2)
        with s1:
            slack_channel = st.text_input("채널명", placeholder="#고객사-채널")
        with s2:
            slack_keyword = st.text_input("고객사명/키워드", placeholder="ABC회사")
        if st.button("Slack 불러오기", type="secondary"):
            with st.spinner("Slack 메시지 수집 중..."):
                parts = []
                if slack_channel:
                    parts.append(f"[채널: {slack_channel}]\n" + fetch_channel_messages(slack_token, slack_channel))
                if slack_keyword:
                    parts.append(f"[검색: {slack_keyword}]\n" + search_messages_by_keyword(slack_token, slack_keyword))
                if parts:
                    st.session_state["slack_data"] = "\n\n".join(parts)
                    st.success(f"✅ {len(st.session_state['slack_data'])}자 수집 완료")
    else:
        st.info("Slack Token 입력 시 채널 대화를 자동으로 가져옵니다")

with tab3:
    w1, w2 = st.columns(2)
    with w1:
        customer_url = st.text_input("고객사 웹사이트", placeholder="https://customer.com")
    with w2:
        industry_keyword = st.text_input("업계 키워드", placeholder="예) 중장년 AI 디지털 교육")
    if st.button("웹 데이터 수집", type="secondary"):
        with st.spinner("웹 데이터 수집 중..."):
            if customer_url:
                st.session_state["customer_web"] = scrape_website(customer_url)
            if industry_keyword:
                st.session_state["industry_news"] = search_industry_news(industry_keyword)
        collected = []
        if st.session_state.get("customer_web"):
            collected.append("고객사 사이트")
        if st.session_state.get("industry_news"):
            collected.append("업계 뉴스")
        if collected:
            st.success(f"✅ {', '.join(collected)} 수집 완료")

st.divider()

# ── 출력 모드 선택 ─────────────────────────────────────────────
st.markdown("### ⚙️ 출력 방식")
output_mode = st.radio(
    "출력 방식",
    options=["🌐 웹 분석 보고서", "📊 PPT 다운로드"],
    horizontal=True,
    label_visibility="collapsed",
)

st.divider()

# ── 소스 요약 + 생성 버튼 ─────────────────────────────────────
sources = []
if rfp_combined:                           sources.append("📋 과업지시서")
if free_text:                              sources.append("📝 메모")
if file_text:                              sources.append(f"📁 파일 {len(uploaded_files)}개")
if st.session_state.get("slack_data"):     sources.append("💬 Slack")
if st.session_state.get("customer_web"):   sources.append("🌐 고객사 웹")
if st.session_state.get("industry_news"):  sources.append("📰 업계 뉴스")

if sources:
    badges = " ".join(f'<span class="source-badge">{s}</span>' for s in sources)
    st.markdown(f"**수집된 소스:** {badges}", unsafe_allow_html=True)
    st.markdown("")

btn_label = "🚀  분석 보고서 생성하기" if output_mode == "🌐 웹 분석 보고서" else "🚀  PPT 생성하기"
generate_btn = st.button(btn_label, type="primary", use_container_width=True)

# ── 생성 실행 ─────────────────────────────────────────────────
if generate_btn:
    if not sources:
        st.error("과업지시서 또는 참고 자료를 하나 이상 입력해주세요")
        st.stop()

    combined = "\n\n".join(filter(None, [
        f"[고객사명]\n{customer_name}" if customer_name else "",
        f"[과업지시서 / RFP - 최우선 참고]\n{rfp_combined}" if rfp_combined else "",
        f"[추가 메모]\n{free_text}" if free_text else "",
        f"[참고 파일]\n{file_text}" if file_text else "",
        st.session_state.get("slack_data", ""),
        st.session_state.get("customer_web", ""),
        st.session_state.get("industry_news", ""),
    ]))

    progress = st.progress(0, text="분석 중...")

    try:
        # ── 웹 분석 보고서 ──────────────────────────────────────
        if output_mode == "🌐 웹 분석 보고서":
            if not anthropic_key:
                progress.empty()
                st.error("웹 분석 보고서는 Anthropic API Key가 필요합니다")
                st.stop()

            from modules.web_analyzer import analyze_for_web
            from modules.web_renderer import render_section_html, render_full_page, section_display_height
            import streamlit.components.v1 as components

            progress.progress(20, text="Claude AI가 자료를 분석 중...")
            analysis = analyze_for_web(combined, anthropic_key)
            progress.progress(80, text="보고서 렌더링 중...")

            st.success(f"✅ 분석 완료 — {analysis.get('title', '')}")

            # 각 섹션을 순서대로 표시
            for section in analysis.get("sections", []):
                section_html = render_section_html(section)
                height = section_display_height(section)
                components.html(section_html, height=height, scrolling=False)

            # 전체 HTML — session_state에 저장해서 리런 후에도 다운로드 가능
            full_html = render_full_page(analysis)
            st.session_state["_web_html"]  = full_html
            st.session_state["_web_fname"] = f"{customer_name or '분석보고서'}_사전분석.html"

        # ── PPT 다운로드 ────────────────────────────────────────
        else:
            if anthropic_key:
                from modules.ai_analyzer import analyze_and_generate_storyline
                progress.progress(20, text="Claude AI가 과업지시서를 분석 중...")
                storyline = analyze_and_generate_storyline(
                    direct_input=combined,
                    api_key=anthropic_key,
                )
            else:
                from modules.template_analyzer import analyze_without_api
                progress.progress(20, text="템플릿 모드로 구조화 중...")
                storyline = analyze_without_api(combined, customer_name)

            progress.progress(70, text="슬라이드 생성 중...")
            pptx_bytes = generate_pptx(storyline)
            progress.progress(100, text="완료!")

            st.success("✅ PPT 생성 완료!")

            with st.expander("📋 생성된 스토리라인", expanded=True):
                st.markdown(f"### {storyline.get('title', '')}")
                st.info(f"**핵심 메시지:** {storyline.get('one_line_summary', '')}")
                st.caption(f"**논리 흐름:** {storyline.get('narrative_thread', '')}")
                st.divider()
                for i, slide in enumerate(storyline.get("slides", []), 1):
                    st.markdown(f"**슬라이드 {i}: {slide.get('title', '')}**")
                    if slide.get("so_what"):
                        st.caption(f"💬 {slide['so_what']}")
                    for pt in slide.get("key_points", []):
                        st.markdown(f"  {pt}")

            st.download_button(
                label="⬇️ PPTX 다운로드",
                data=pptx_bytes,
                file_name=f"{customer_name or '제안서'}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                type="primary",
                use_container_width=True,
            )

    except Exception as e:
        progress.empty()
        st.error(f"생성 실패: {str(e)}")
        st.exception(e)

# ── 웹 보고서 다운로드 버튼 (리런 후에도 유지) ─────────────────
if "_web_html" in st.session_state and output_mode == "🌐 웹 분석 보고서":
    st.download_button(
        label="⬇️ HTML 보고서 다운로드",
        data=st.session_state["_web_html"].encode("utf-8"),
        file_name=st.session_state["_web_fname"],
        mime="text/html",
        type="primary",
        use_container_width=True,
    )
