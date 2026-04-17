import io
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

COLOR_PRIMARY = RGBColor(0x1A, 0x56, 0xDB)
COLOR_ACCENT  = RGBColor(0x0E, 0x9F, 0x6E)
COLOR_RED     = RGBColor(0xD6, 0x17, 0x1A)
COLOR_DARK    = RGBColor(0x11, 0x18, 0x27)
COLOR_LIGHT   = RGBColor(0xF9, 0xFA, 0xFB)
COLOR_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_GRAY    = RGBColor(0x6B, 0x72, 0x80)
COLOR_BGLINE  = RGBColor(0xE5, 0xE7, 0xEB)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

SECTION_COLORS = {
    "background":  COLOR_PRIMARY,
    "issues":      COLOR_RED,
    "problems":    COLOR_RED,
    "solution":    RGBColor(0x05, 0x7A, 0x55),
    "impact":      RGBColor(0x17, 0x71, 0xD3),
    "toc":         COLOR_DARK,
    "rfp":         RGBColor(0xD9, 0x77, 0x06),
    "credentials": RGBColor(0x05, 0x7A, 0x55),
    "timeline":    RGBColor(0x76, 0x54, 0xF8),
}
SECTION_ICONS = {
    "background":  "📋",
    "issues":      "⚠️",
    "problems":    "🔍",
    "solution":    "💡",
    "impact":      "🎯",
    "toc":         "📑",
    "rfp":         "📌",
    "credentials": "🏆",
    "timeline":    "📅",
}


def _bg(slide, color=COLOR_LIGHT):
    s = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()


def _rect(slide, x, y, w, h, fill, line_color=None, line_w=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color
        if line_w: s.line.width = line_w
    else:
        s.line.fill.background()
    return s


def _txt(slide, text, x, y, w, h,
         size=11, bold=False, color=COLOR_DARK,
         align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def _add_cover(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, Inches(5.8), SLIDE_H, COLOR_PRIMARY)
    _rect(slide, Inches(5.8), 0, Inches(7.53), SLIDE_H, COLOR_LIGHT)
    _rect(slide, 0, 0, Inches(5.8), Inches(0.08), RGBColor(0x0E, 0x9F, 0x6E))

    _txt(slide, "BUSINESS PROPOSAL",
         Inches(0.5), Inches(1.2), Inches(5.0), Inches(0.5),
         size=9, bold=False, color=RGBColor(0x93, 0xC5, 0xFD), italic=True)

    _txt(slide, data.get("title", "제안서"),
         Inches(0.5), Inches(1.8), Inches(5.0), Inches(2.8),
         size=26, bold=True, color=COLOR_WHITE, wrap=True)

    _txt(slide, data.get("subtitle", ""),
         Inches(0.5), Inches(4.8), Inches(5.0), Inches(1.0),
         size=13, color=RGBColor(0xBF, 0xDB, 0xFF), wrap=True)

    _txt(slide, "데이원컴퍼니 (패스트캠퍼스)",
         Inches(0.5), Inches(6.5), Inches(5.0), Inches(0.6),
         size=10, color=RGBColor(0x93, 0xC5, 0xFD))


def _add_content(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, COLOR_LIGHT)

    stype = data.get("type", "background")
    accent = SECTION_COLORS.get(stype, COLOR_PRIMARY)
    icon   = SECTION_ICONS.get(stype, "")

    # 상단 바
    _rect(slide, 0, 0, SLIDE_W, Inches(0.06), accent)

    y = Inches(0.1)

    # 제목 영역
    title_text = f"{icon}  {data.get('title','')}" if icon else data.get("title","")
    _txt(slide, title_text,
         Inches(0.4), y, Inches(12.5), Inches(0.58),
         size=18, bold=True, color=COLOR_DARK)
    y += Inches(0.6)

    # 부제목 (1. 섹션명 > 세부주제)
    subtitle = data.get("subtitle", "")
    if subtitle:
        _txt(slide, subtitle,
             Inches(0.4), y, Inches(12.5), Inches(0.24),
             size=9, color=COLOR_GRAY, italic=True)
        y += Inches(0.24)

    # So-What 박스
    so_what = data.get("so_what", "")
    if so_what:
        _rect(slide, Inches(0.4), y, Inches(12.5), Inches(0.50), accent)
        _txt(slide, f"▶  {so_what}",
             Inches(0.55), y + Inches(0.04), Inches(12.2), Inches(0.46),
             size=10, bold=True, color=COLOR_WHITE, wrap=True)
        y += Inches(0.58)

    body_top = y
    body_h   = SLIDE_H - body_top - Inches(0.15)

    key_points = data.get("key_points", [])
    details    = data.get("details", "")
    sub_points = data.get("sub_points", [])  # 선택적 세부 포인트

    if not key_points:
        return

    # ── 레이아웃 결정 ────────────────────────────────────────
    # 포인트 수에 따라 1열 or 2열
    use_two_col = len(key_points) >= 5 or any(len(p) > 80 for p in key_points)

    if use_two_col:
        _render_two_col(slide, key_points, sub_points, details,
                        body_top, body_h, accent)
    else:
        _render_single_col(slide, key_points, sub_points, details,
                           body_top, body_h, accent)


def _render_single_col(slide, points, sub_points, details,
                       top, height, accent):
    """포인트 4개 이하 — 1열 넓게"""
    n = len(points)
    item_h = (height - Inches(0.3)) / max(n, 1)
    item_h = min(item_h, Inches(1.6))

    for i, pt in enumerate(points):
        y = top + i * item_h

        # 번호 원
        _rect(slide, Inches(0.4), y + Inches(0.08),
              Inches(0.38), Inches(0.38), accent)
        _txt(slide, str(i+1),
             Inches(0.4), y + Inches(0.08), Inches(0.38), Inches(0.38),
             size=10, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

        # 메인 포인트 텍스트
        _txt(slide, pt,
             Inches(0.95), y, Inches(11.9), item_h - Inches(0.05),
             size=12, color=COLOR_DARK, wrap=True)

    # 하단 상세
    if details:
        det_y = top + n * item_h + Inches(0.05)
        _rect(slide, Inches(0.4), det_y, Inches(12.5), Inches(0.02),
              RGBColor(0xD1, 0xD5, 0xDB))
        _txt(slide, f"📌  {details}",
             Inches(0.4), det_y + Inches(0.08), Inches(12.5), Inches(0.9),
             size=9, color=COLOR_GRAY, italic=True, wrap=True)


def _render_two_col(slide, points, sub_points, details,
                    top, height, accent):
    """포인트 5개+ — 2열 레이아웃, 텍스트 최대 표시"""
    half  = (len(points) + 1) // 2
    col_w = Inches(6.3)
    gap   = Inches(0.4)
    item_h = (height - Inches(0.4)) / max(half, 1)
    item_h = min(item_h, Inches(1.4))

    for i, pt in enumerate(points):
        col = i // half
        row = i % half
        x   = Inches(0.4) + col * (col_w + gap)
        y   = top + row * item_h

        # 번호 뱃지
        _rect(slide, x, y + Inches(0.06), Inches(0.32), Inches(0.32), accent)
        _txt(slide, str(i+1),
             x, y + Inches(0.06), Inches(0.32), Inches(0.32),
             size=9, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

        # 포인트 텍스트 — 남은 공간 전부 사용
        _txt(slide, pt,
             x + Inches(0.4), y, col_w - Inches(0.45), item_h - Inches(0.05),
             size=10, color=COLOR_DARK, wrap=True)

    # 상세 설명 (하단 전체 폭)
    if details:
        det_y = top + half * item_h + Inches(0.05)
        if det_y + Inches(0.6) < SLIDE_H:
            _rect(slide, Inches(0.4), det_y, Inches(12.5), Inches(0.02),
                  RGBColor(0xD1, 0xD5, 0xDB))
            _txt(slide, f"📌  {details}",
                 Inches(0.4), det_y + Inches(0.06), Inches(12.5),
                 SLIDE_H - det_y - Inches(0.1),
                 size=9, color=COLOR_GRAY, italic=True, wrap=True)


def _generate_legacy(storyline: dict) -> bytes:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    for sd in storyline.get("slides", []):
        if sd.get("type") == "cover":
            _add_cover(prs, sd)
        else:
            _add_content(prs, sd)
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _generate_playwright(storyline: dict) -> bytes:
    from playwright.sync_api import sync_playwright
    from modules.html_renderer import render_slide

    slides_data = storyline.get("slides", [])
    total = len(slides_data)
    images = []

    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": 1280, "height": 720})
        for i, sd in enumerate(slides_data, 1):
            html = render_slide(sd, total, i)
            page.set_content(html)
            page.wait_for_timeout(180)
            images.append(page.screenshot(type="png"))
        browser.close()

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    for img_bytes in images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(
            io.BytesIO(img_bytes), 0, 0,
            prs.slide_width, prs.slide_height,
        )
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


def generate_pptx(storyline: dict) -> bytes:
    try:
        return _generate_playwright(storyline)
    except Exception:
        return _generate_legacy(storyline)
